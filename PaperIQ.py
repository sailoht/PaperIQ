import streamlit as st
import os
import pandas as pd
import json
import csv
import docx2txt
import pytesseract
from PIL import Image
from pptx import Presentation
from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain.embeddings import OpenAIEmbeddings
from langchain.llms import OpenAI
from langchain.chains import ConversationalRetrievalChain

# ------------------------
# APP CONFIG
# ------------------------
st.set_page_config(page_title="PaperIQ", layout="wide")
st.title("📘 PaperIQ – Smart Document Assistant")
st.caption("Built by Sai Lohith | Powered by LangChain & PrivateGPT base")

# ------------------------
# SIDEBAR CONFIG
# ------------------------
with st.sidebar:
    st.header("Configuration")
    api_key = st.text_input("Enter your OpenAI API Key", type="password")
    if not api_key:
        st.warning("Please enter your OpenAI API key to continue.")
    uploaded_files = st.file_uploader(
        "Upload your documents",
        type=["pdf", "txt", "docx", "pptx", "xlsx", "csv", "json", "md", "jpg", "jpeg", "png"],
        accept_multiple_files=True
    )

# ------------------------
# TEXT EXTRACTION FUNCTION
# ------------------------
def extract_text_from_file(file):
    ext = file.name.lower().split(".")[-1]
    temp_path = os.path.join("data", file.name)
    with open(temp_path, "wb") as f:
        f.write(file.getbuffer())

    if ext == "pdf":
        loader = PyPDFLoader(temp_path)
        docs = loader.load()
        return " ".join([d.page_content for d in docs])

    elif ext == "txt":
        loader = TextLoader(temp_path)
        docs = loader.load()
        return " ".join([d.page_content for d in docs])

    elif ext == "docx":
        return docx2txt.process(temp_path)

    elif ext == "pptx":
        prs = Presentation(temp_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    elif ext in ["xlsx", "csv"]:
        try:
            if ext == "csv":
                df = pd.read_csv(temp_path)
            else:
                df = pd.read_excel(temp_path)
            return df.to_string()
        except Exception as e:
            return f"Error reading table: {e}"

    elif ext == "json":
        with open(temp_path, "r", encoding="utf-8") as f:
            return json.dumps(json.load(f), indent=2)

    elif ext == "md":
        with open(temp_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext in ["jpg", "jpeg", "png"]:
        image = Image.open(temp_path)
        text = pytesseract.image_to_string(image)
        return text if text.strip() else "[No readable text found in image]"

    else:
        return f"Unsupported file format: {ext}"

# ------------------------
# CORE LOGIC
# ------------------------
if api_key and uploaded_files:
    os.environ["OPENAI_API_KEY"] = api_key

    all_texts = []
    for file in uploaded_files:
        st.write(f"Processing: {file.name}")
        extracted_text = extract_text_from_file(file)
        if extracted_text:
            all_texts.append(extracted_text)

    if not all_texts:
        st.error(" No readable content found in uploaded files.")
        st.stop()

    combined_text = "\n\n".join(all_texts)
    st.success(f"Extracted text from {len(uploaded_files)} files.")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
    split_docs = text_splitter.create_documents([combined_text])
    embeddings = OpenAIEmbeddings()
    vector_store = Chroma.from_documents(split_docs, embeddings)

    retriever = vector_store.as_retriever()
    llm = OpenAI(temperature=0)
    qa_chain = ConversationalRetrievalChain.from_llm(llm, retriever=retriever)

    chat_history = []
    st.subheader("Chat with your documents")

    user_input = st.text_input("Ask a question:")
    if user_input:
        response = qa_chain({"question": user_input, "chat_history": chat_history})
        st.markdown(f"**Answer:** {response['answer']}")
        chat_history.append((user_input, response['answer']))

    if st.button(" Summarize All Documents"):
        summary_prompt = "Summarize all these documents briefly and clearly."
        response = qa_chain({"question": summary_prompt, "77chat_history": []})
        st.subheader("Summary")
        st.write(response['answer'])

else:
    st.info(" Upload your documents and enter your API key to begin.")
